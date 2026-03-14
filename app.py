import streamlit as st
import requests
import json
import graphviz
import io
from pptx import Presentation
from pptx.util import Pt

# 1. Flowchart Logic (Requirement #11: logic to structured steps)
def generate_local_flowchart(steps):
    """Renders a visual flowchart from AI-provided steps locally."""
    dot = graphviz.Digraph(comment='Process Flowchart')
    dot.attr(bgcolor='#0f172a', color='white', fontcolor='white')
    dot.attr('node', shape='box', style='rounded,filled', 
             fillcolor='#1e293b', color='#6366f1', fontcolor='white')
    
    for i, step in enumerate(steps):
        dot.node(str(i), step)
        if i > 0:
            dot.edge(str(i-1), str(i), color='#a855f7')
    return dot

# 2. PPTX Generation Logic (Requirement #7 & #12: Export/Renderer)
def create_pptx_file(presentation_data):
    prs = Presentation()
    for slide_data in presentation_data['slides']:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        slide.shapes.title.text = slide_data['title']
        
        # Set Content
        body = slide.placeholders[1]
        body.text = slide_data['content']
        
    binary_output = io.BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

# 3. Local AI Engine (Requirement #7: Support 2-3 page mode)
def call_local_ai(prompt, mode):
    slide_count = 3 if mode == "Short (2-3 pages)" else 6
    url = "http://localhost:11434/api/generate"
    
    # Prompting for structured JSON including flowchart data
    full_prompt = f"""
    Create a {slide_count} slide presentation outline for: {prompt}.
    Include slide titles, content, and a 'flowchart_steps' list of 3-5 logical steps for the topic.
    Return ONLY JSON format:
    {{
      "slides": [{{ "title": "Slide Title", "content": "Bullet points here" }}],
      "flowchart_steps": ["Step A", "Step B", "Step C"]
    }}
    """
    
    try:
        response = requests.post(url, json={
            "model": "llama3",
            "prompt": full_prompt,
            "stream": False,
            "format": "json"
        }, timeout=60)
        return json.loads(response.json()['response'])
    except Exception as e:
        st.error(f"Connection Error: Ensure Ollama is running. {e}")
        return None

# 4. Interactive UI (Requirement #8: User-friendly interface)
def main():
    st.set_page_config(page_title="Gamma Local AI", layout="wide")
    
    # Custom CSS for Gamma-style Dark Mode
    st.markdown("""
        <style>
        .stApp { background-color: #0f172a; color: white; }
        .slide-card { 
            background: rgba(255, 255, 255, 0.05); 
            padding: 20px; border-radius: 15px; 
            border: 1px solid #6366f1; margin-bottom: 20px;
        }
        </style>
    """, unsafe_allow_html=True)

    st.title("✨ Gamma-Style Local AI")
    st.subheader("Automated presentations and flowcharts (No API keys required)")

    with st.sidebar:
        st.header("Settings")
        mode = st.radio("Presentation Length", ["Standard", "Short (2-3 pages)"]) # Requirement #86
        topic = st.text_input("Topic", placeholder="e.g. How Solar Panels Work")
        generate_btn = st.button("Generate Everything")

    if generate_btn and topic:
        with st.spinner("AI is thinking and drawing diagrams..."):
            result = call_local_ai(topic, mode)
            if result:
                st.session_state.presentation = result

    if "presentation" in st.session_state:
        data = st.session_state.presentation
        
        # Display Slides
        for slide in data['slides']:
            st.markdown(f"<div class='slide-card'><h3>{slide['title']}</h3><p>{slide['content']}</p></div>", unsafe_allow_html=True)

        # Display Flowchart (Requirement #103, #104)
        if data.get('flowchart_steps'):
            st.divider()
            st.subheader("📊 Logic Flowchart")
            chart = generate_local_flowchart(data['flowchart_steps'])
            st.graphviz_chart(chart)

        # Export (Requirement #76)
        st.divider()
        ppt_file = create_pptx_file(data)
        st.download_button(
            label="💾 Download for Microsoft PowerPoint",
            data=ppt_file,
            file_name="Gamma_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()