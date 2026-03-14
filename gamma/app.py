"""
Gamma-Local v1.0 — AI-Powered Presentation Generator (100% Local)
==================================================================
Implements all Functional Requirements and Technical Constraints
from Gamma-Local PRD v1.0.

Run:
    streamlit run app.py

Requirements:
    pip install -r requirements.txt
    ollama serve   (in a separate terminal)
    ollama pull llama3:8b
"""

# ─────────────────────────────────────────────────────────────────────────────
# Standard library
# ─────────────────────────────────────────────────────────────────────────────
import io
import json
import os
import re
import shutil
import tempfile
import time
import unicodedata
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any

# ─────────────────────────────────────────────────────────────────────────────
# Third-party — all local, zero data egress (PRD §4.4)
# ─────────────────────────────────────────────────────────────────────────────
import jsonschema
import ollama
import streamlit as st
from docx import Document as DocxDocument
from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# PDF extraction: pdfminer.six primary, PyMuPDF fallback (PRD §3.1.2)
try:
    from pdfminer.high_level import extract_text as pdfminer_extract
    HAS_PDFMINER = True
except ImportError:
    HAS_PDFMINER = False

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

# Token counting (PRD §3.2.3)
try:
    import tiktoken
    _enc = tiktoken.get_encoding("cl100k_base")
    def count_tokens(text: str) -> int:
        return len(_enc.encode(text))
    HAS_TIKTOKEN = True
except Exception:
    def count_tokens(text: str) -> int:        # fallback: ~0.75 words/token
        return int(len(text.split()) / 0.75)
    HAS_TIKTOKEN = False

# ─────────────────────────────────────────────────────────────────────────────
# Configuration (PRD §Appendix B)
# ─────────────────────────────────────────────────────────────────────────────
CFG = {
    "ollama": {
        "host":            "http://localhost:11434",
        "default_model":   "llama3:8b",
        "request_timeout": 120,
        "supported_models": ["llama3:8b", "mistral:7b", "llama3:70b", "phi3:mini"],
    },
    "ingestion": {
        "max_file_size_mb": 50,
        "chunk_size_tokens": 800,
    },
    "generation": {
        "max_slides":             30,
        "min_slides":              8,
        "retry_attempts":          3,
        "parallel_refinement": False,
        "map_reduce_threshold":  4000,   # tokens; PRD §3.2.3
    },
    "output": {
        "directory":          str(Path.home() / "gamma-local-output"),
        "cleanup_temp_files": True,
    },
    "security": {
        "enable_audit_log": False,
    },
}

OUTPUT_DIR = Path(CFG["output"]["directory"])
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Loguru — structured local-only logging (PRD §4.4, no remote sink)
logger.add(
    OUTPUT_DIR / "gamma_local.log",
    rotation="10 MB",
    retention="30 days",
    level="DEBUG",
    format="{time:YYYY-MM-DD HH:mm:ss} | {level} | {message}",
)

# ─────────────────────────────────────────────────────────────────────────────
# Design tokens — "Ocean Gradient" palette (professional, tech-appropriate)
# ─────────────────────────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x06, 0x5A, 0x82)   # deep blue  — dominant
C_TEAL    = RGBColor(0x1C, 0x72, 0x93)   # teal       — secondary
C_MIDNIGHT= RGBColor(0x21, 0x29, 0x5C)   # midnight   — dark bg
C_ICE     = RGBColor(0xD6, 0xEA, 0xF8)   # ice blue   — card bg
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TXT= RGBColor(0x1A, 0x25, 0x2F)   # near-black body text
C_ACCENT  = RGBColor(0x02, 0xC3, 0x9A)   # mint green — accent / bullets

SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# JSON schema for Stage-1 LLM output validation (PRD §3.2.4)
# ─────────────────────────────────────────────────────────────────────────────
SLIDE_SCHEMA = {
    "type": "array",
    "minItems": 1,
    "items": {
        "type": "object",
        # Only slide_number and layout required; smaller models omit content on some slides
        "required": ["slide_number", "layout"],
        "properties": {
            "slide_number": {"type": ["integer", "string"]},
            "layout": {"type": "string"},
            "title":   {"type": "string"},
            "content": {},
        },
    },
}

# Stage-2 refined slide schemas
REFINED_BULLETS_SCHEMA = {
    "type": "object",
    "required": ["title", "bullets"],
    "properties": {
        "title":   {"type": "string"},
        "bullets": {"type": "array", "items": {"type": "string"}, "maxItems": 6},
    },
}
REFINED_TWOCOL_SCHEMA = {
    "type": "object",
    "required": ["title", "left", "right"],
    "properties": {
        "title": {"type": "string"},
        "left":  {"type": "string"},
        "right": {"type": "string"},
    },
}

# ─────────────────────────────────────────────────────────────────────────────
# §4.3 — Ollama health-check and environment verification
# ─────────────────────────────────────────────────────────────────────────────

def check_ollama_binary() -> bool:
    """Verify ollama binary is on PATH (PRD §4.3 check 1)."""
    return shutil.which("ollama") is not None


def check_ollama_server() -> tuple[bool, list[str]]:
    """
    PRD §4.3 checks 2 & 3:
    - Server reachable at localhost:11434
    - At least one supported model is available
    Returns (ok: bool, available_models: list[str])
    """
    try:
        client = ollama.Client(host=CFG["ollama"]["host"])
        resp   = client.list()
        # ollama library may return models as objects or dicts depending on version
        names  = []
        for m in resp.models:
            try:
                name = m.model  # newer ollama-python
            except AttributeError:
                try:
                    name = m["model"]  # dict fallback
                except (KeyError, TypeError):
                    name = str(m)
            if name:
                names.append(name)
        return True, names
    except Exception as e:
        logger.warning(f"Ollama server check failed: {e}")
        return False, []


def get_best_available_model(available: list[str]) -> str | None:
    """Return highest-priority supported model from available list."""
    SUPPORTED_BASES = ["llama3", "mistral", "phi3", "llama2", "gemma"]
    priority = CFG["ollama"]["supported_models"]  # e.g. ["llama3:8b", ...]

    # 1. Exact match first
    for preferred in priority:
        if preferred in available:
            return preferred

    # 2. Base-name prefix match (handles llama3:latest, llama3:8b-instruct, etc.)
    for preferred in priority:
        base = preferred.split(":")[0]
        for avail in available:
            if avail.split(":")[0] == base:
                return avail

    # 3. Any known-good base
    for base in SUPPORTED_BASES:
        for avail in available:
            if avail.split(":")[0] == base:
                return avail

    # 4. Whatever is available
    return available[0] if available else None


# ─────────────────────────────────────────────────────────────────────────────
# §3.1 — File ingestion pipeline
# ─────────────────────────────────────────────────────────────────────────────

def validate_file(file_bytes: bytes, filename: str) -> None:
    """PRD §3.1.2 — validate extension and size."""
    ext  = Path(filename).suffix.lower()
    allowed = {".pdf", ".txt", ".md", ".docx"}
    if ext not in allowed:
        raise ValueError(
            f"Unsupported format '{ext}'. Supported: {', '.join(allowed)}"
        )
    max_mb = CFG["ingestion"]["max_file_size_mb"]
    size_mb = len(file_bytes) / (1024 * 1024)
    if size_mb > max_mb:
        raise ValueError(
            f"File too large ({size_mb:.1f} MB). Maximum: {max_mb} MB."
        )


def extract_pdf_text(raw: bytes) -> str:
    """pdfminer.six primary, PyMuPDF fallback (PRD §3.1.2)."""
    if HAS_PDFMINER:
        try:
            return pdfminer_extract(io.BytesIO(raw))
        except Exception as e:
            logger.warning(f"pdfminer failed ({e}), trying PyMuPDF...")

    if HAS_PYMUPDF:
        try:
            doc   = fitz.open(stream=raw, filetype="pdf")
            pages = [page.get_text() for page in doc]
            return "\n\n".join(pages)
        except fitz.FileDataError:
            raise RuntimeError(
                "PDF appears to be encrypted or corrupt. "
                "Please remove password protection before uploading."
            )

    raise RuntimeError(
        "No PDF extraction library found. "
        "Run: pip install pdfminer.six"
    )


def extract_text(file_bytes: bytes, filename: str) -> str:
    """Dispatch extraction by file type (PRD §3.1.2)."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return extract_pdf_text(file_bytes)

    if ext in (".txt", ".md"):
        return file_bytes.decode("utf-8", errors="replace")

    if ext == ".docx":
        doc   = DocxDocument(io.BytesIO(file_bytes))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paras)

    raise ValueError(f"Unsupported format: {ext}")


def clean_text(raw: str) -> str:
    """
    PRD §3.1.2 — Text Cleaning & Normalization:
    - Unicode NFC normalisation
    - Remove control characters (except newlines/tabs)
    - Collapse excessive whitespace
    - De-hyphenate line-wrapped words
    - Strip obvious page number / header patterns
    """
    text = unicodedata.normalize("NFC", raw)

    # Remove control characters (keep \n \t)
    text = re.sub(r"[^\S\n\t ]+", " ", text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)

    # De-hyphenate: "compli-\ncated" → "complicated"
    text = re.sub(r"-\s*\n\s*", "", text)

    # Remove obvious page-number lines (lines that are only digits / "Page N of M")
    text = re.sub(r"(?m)^\s*(Page\s+\d+\s+(of\s+\d+)?|\d+)\s*$", "", text)

    # Collapse 3+ blank lines → 2
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def chunk_text(text: str, max_tokens: int = None) -> list[str]:
    """
    PRD §3.1.2 — Chunking & Sectioning.
    Split on detected section boundaries, then hard-cap at max_tokens.
    """
    max_tokens = max_tokens or CFG["ingestion"]["chunk_size_tokens"]

    # Section boundary patterns
    section_re = re.compile(
        r"(?m)^(#{1,3}\s.+|"                    # Markdown headings
        r"(?:Introduction|Background|Summary|"
        r"Requirements?|Objective|Conclusion|"
        r"Appendix|Overview|Goals?|Scope)\s*:?)"
        r"\s*$",
        re.IGNORECASE,
    )

    # Split into candidate chunks at headings
    boundaries = [m.start() for m in section_re.finditer(text)]
    if not boundaries:
        boundaries = []
    boundaries = [0] + boundaries + [len(text)]

    raw_chunks: list[str] = []
    for i in range(len(boundaries) - 1):
        chunk = text[boundaries[i]:boundaries[i + 1]].strip()
        if chunk:
            raw_chunks.append(chunk)

    # Hard-cap each chunk at max_tokens
    final_chunks: list[str] = []
    for chunk in raw_chunks:
        if count_tokens(chunk) <= max_tokens:
            final_chunks.append(chunk)
        else:
            # Split by sentences into sub-chunks
            sentences = re.split(r"(?<=[.!?])\s+", chunk)
            buf, buf_tok = [], 0
            for sent in sentences:
                t = count_tokens(sent)
                if buf_tok + t > max_tokens and buf:
                    final_chunks.append(" ".join(buf))
                    buf, buf_tok = [sent], t
                else:
                    buf.append(sent)
                    buf_tok += t
            if buf:
                final_chunks.append(" ".join(buf))

    return final_chunks if final_chunks else [text]


def validate_content_length(text: str) -> None:
    """PRD §3.1.3 — halt on insufficient content."""
    tokens = count_tokens(text)
    if tokens < 100:
        raise ValueError(
            f"Document contains fewer than 100 tokens ({tokens} found). "
            "Please upload a document with sufficient textual content."
        )


# ─────────────────────────────────────────────────────────────────────────────
# §3.2 — AI Orchestration
# ─────────────────────────────────────────────────────────────────────────────

STAGE1_SYSTEM = (
    "You are a professional presentation architect. "
    "Analyze the provided document and produce a structured JSON outline "
    "for a PowerPoint presentation. "
    'The JSON must conform exactly to the following schema: '
    '[{"slide_number": int, "layout": "TITLE" | "SECTION_HEADER" | "BULLETS" | '
    '"TWO_COLUMN" | "QUOTE" | "CLOSING", "title": string, "content": string | '
    '[string] | {"left": string, "right": string}}]. '
    f'Generate between {CFG["generation"]["min_slides"]} and '
    f'{CFG["generation"]["max_slides"]} slides. '
    "Prioritize clarity, logical flow, and executive-level communication. "
    "Return ONLY valid JSON. "
    "Do not include any explanation, preamble, or markdown formatting outside the JSON array."
)

STAGE2_SYSTEM = (
    "You are an expert presentation writer. "
    "Rewrite the following slide content into concise, professional presentation language. "
    "Slide title must be 8 words or fewer. "
    "Each bullet point must be 20 words or fewer. "
    "Maintain parallel grammatical structure. Use active voice. "
    'Return ONLY a JSON object: {"title": string, "bullets": [string]} '
    'or {"title": string, "left": string, "right": string} as appropriate for the layout.'
)

STAGE1_MAP_SYSTEM = (
    "You are a professional document analyst. "
    "Summarize the following document section into 3–5 concise sentences "
    "capturing the key points, decisions, and requirements. "
    "Return plain text only."
)


def _call_ollama(model: str, system: str, user: str,
                 timeout: int = None) -> str:
    """
    Call local Ollama via the ollama Python library.
    All traffic stays on 127.0.0.1 (PRD §4.4).
    """
    timeout = timeout or CFG["ollama"]["request_timeout"]
    client  = ollama.Client(
        host=CFG["ollama"]["host"],
    )
    resp = client.chat(
        model=model,
        messages=[
            {"role": "system",  "content": system},
            {"role": "user",    "content": user},
        ],
        options={"temperature": 0.3, "num_predict": 4096},
    )
    return resp.message.content.strip()


def _strip_json_fences(raw: str) -> str:
    """Strip ```json … ``` fences if the model wraps its output."""
    raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.MULTILINE)
    raw = re.sub(r"\s*```$", "",      raw, flags=re.MULTILINE)
    return raw.strip()


def _parse_json_safe(raw: str) -> Any:
    """
    Find and parse the first valid JSON value in a string.
    Also normalises slide objects so missing fields don't cause schema failures.
    """
    raw = _strip_json_fences(raw)

    parsed = None
    # Try direct parse first
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        pass

    if parsed is None:
        # Find the first [ or { and attempt to parse from there
        for opener, closer in [("[", "]"), ("{", "}")]:
            start = raw.find(opener)
            if start == -1:
                continue
            end = raw.rfind(closer)
            if end == -1:
                continue
            try:
                parsed = json.loads(raw[start:end + 1])
                break
            except json.JSONDecodeError:
                continue

    if parsed is None:
        raise ValueError(f"No valid JSON found in model output:\n{raw[:500]}")

    # Normalise: if model returned a dict with a 'slides' key, unwrap it
    if isinstance(parsed, dict):
        for key in ("slides", "outline", "presentation"):
            if key in parsed and isinstance(parsed[key], list):
                parsed = parsed[key]
                break
        else:
            parsed = [parsed]

    # Normalise each slide object
    VALID_LAYOUTS = {"TITLE", "SECTION_HEADER", "BULLETS", "TWO_COLUMN", "QUOTE", "CLOSING"}
    for i, slide in enumerate(parsed):
        if not isinstance(slide, dict):
            continue
        # Ensure slide_number is int
        slide["slide_number"] = int(slide.get("slide_number", i + 1))
        # Uppercase and validate layout
        layout = str(slide.get("layout", "BULLETS")).upper().strip()
        if layout not in VALID_LAYOUTS:
            layout = "BULLETS"
        slide["layout"] = layout
        # Ensure title exists
        if "title" not in slide or not slide["title"]:
            slide["title"] = f"Slide {slide['slide_number']}"
        # Ensure content exists
        if "content" not in slide:
            slide["content"] = ""

    return parsed


def map_reduce_summarize(chunks: list[str], model: str) -> str:
    """
    PRD §3.2.3 — Map-Reduce for documents > 4 000 tokens.
    Summarize each chunk independently, then concatenate.
    """
    summaries = []
    for i, chunk in enumerate(chunks):
        logger.info(f"Map-reduce: summarising chunk {i+1}/{len(chunks)}")
        summary = _call_ollama(model, STAGE1_MAP_SYSTEM, chunk)
        summaries.append(summary)
    return "\n\n".join(summaries)


def stage1_generate_outline(text: str, model: str,
                             progress_cb=None) -> list[dict]:
    """
    PRD §3.2.2 Stage 1 — Structural Analysis.
    Includes retry logic (PRD §3.2.4) and map-reduce for long docs (§3.2.3).
    """
    # Map-reduce if document too large
    token_count = count_tokens(text)
    threshold   = CFG["generation"]["map_reduce_threshold"]
    if token_count > threshold:
        if progress_cb:
            progress_cb(f"Document is large ({token_count} tokens). "
                        f"Summarising sections first…")
        chunks = chunk_text(text, CFG["ingestion"]["chunk_size_tokens"])
        text   = map_reduce_summarize(chunks, model)

    retries = CFG["generation"]["retry_attempts"]
    last_err = ""
    for attempt in range(1, retries + 1):
        if progress_cb and attempt > 1:
            progress_cb(f"Retrying Stage 1 (attempt {attempt}/{retries})… "
                        f"Last error: {last_err[:80]}")
        extra = (
            ""
            if attempt == 1
            else (
                f"\n\nPREVIOUS ATTEMPT FAILED WITH: {last_err}\n"
                "You MUST return ONLY a valid JSON array matching the schema. "
                "No preamble, no markdown, no explanation."
            )
        )
        try:
            raw    = _call_ollama(model, STAGE1_SYSTEM, text + extra)
            parsed = _parse_json_safe(raw)
            if not isinstance(parsed, list):
                raise ValueError("Top-level value must be a JSON array")
            jsonschema.validate(instance=parsed, schema=SLIDE_SCHEMA)
            logger.info(f"Stage 1 succeeded on attempt {attempt}: "
                        f"{len(parsed)} slides")
            return parsed
        except (ValueError, jsonschema.ValidationError, json.JSONDecodeError) as e:
            last_err = str(e)
            logger.warning(f"Stage 1 attempt {attempt} failed: {last_err}")

    # Fallback — template-based outline from raw text (PRD §3.2.4)
    logger.error("Stage 1: all retries exhausted, using template fallback")
    st.warning(
        "⚠️ AI structural analysis was unavailable after 3 attempts. "
        "Falling back to template-based generation. "
        "Output quality may be reduced."
    )
    return _template_fallback_outline(text)


def stage2_refine_slide(slide: dict, model: str) -> dict:
    """
    PRD §3.2.2 Stage 2 — Content Refinement per slide.
    Includes retry logic (PRD §3.2.4).
    """
    layout  = slide.get("layout", "BULLETS")
    content = slide.get("content", "")
    user_msg = (
        f"Layout: {layout}\n"
        f"Title: {slide.get('title', '')}\n"
        f"Content: {json.dumps(content, ensure_ascii=False)}\n\n"
        "Rewrite into polished presentation language per instructions."
    )

    retries  = CFG["generation"]["retry_attempts"]
    last_err = ""
    for attempt in range(1, retries + 1):
        extra = (
            "" if attempt == 1
            else f"\n\nERROR IN PREVIOUS ATTEMPT: {last_err}. Return ONLY valid JSON."
        )
        try:
            raw    = _call_ollama(model, STAGE2_SYSTEM, user_msg + extra)
            parsed = _parse_json_safe(raw)
            if not isinstance(parsed, dict):
                raise ValueError("Stage 2 must return a JSON object")

            # Validate schema based on layout
            if layout == "TWO_COLUMN":
                jsonschema.validate(parsed, REFINED_TWOCOL_SCHEMA)
            else:
                jsonschema.validate(parsed, REFINED_BULLETS_SCHEMA)

            return parsed
        except (ValueError, jsonschema.ValidationError, json.JSONDecodeError) as e:
            last_err = str(e)
            logger.warning(f"Stage 2 attempt {attempt} failed for "
                           f"slide '{slide.get('title','')}': {last_err}")

    # Fallback: return raw slide content as-is
    logger.error(f"Stage 2 fallback for slide: {slide.get('title','')}")
    return _stage2_fallback(slide)


def _stage2_fallback(slide: dict) -> dict:
    """Template fallback when Stage 2 AI refinement fails (PRD §3.2.4)."""
    layout  = slide.get("layout", "BULLETS")
    content = slide.get("content", "")
    title   = slide.get("title", "Slide")[:60]

    if layout == "TWO_COLUMN" and isinstance(content, dict):
        return {
            "title": title,
            "left":  str(content.get("left", "")),
            "right": str(content.get("right", "")),
        }

    bullets: list[str] = []
    if isinstance(content, list):
        bullets = [str(b)[:100] for b in content[:6]]
    elif isinstance(content, str):
        # Split paragraph into bullets at sentence boundaries
        sentences = re.split(r"(?<=[.!?])\s+", content)
        bullets   = [s.strip() for s in sentences if s.strip()][:6]

    return {"title": title, "bullets": bullets}


def _template_fallback_outline(text: str) -> list[dict]:
    """Build a minimal outline from raw text when Stage 1 fails (PRD §3.2.4)."""
    chunks  = chunk_text(text, 400)
    outline = [{"slide_number": 1, "layout": "TITLE",
                "title": "Presentation", "content": "Generated by Gamma-Local"}]
    for i, chunk in enumerate(chunks[:14], start=2):
        first_line = chunk.split("\n")[0][:60].strip() or f"Section {i-1}"
        sentences  = re.split(r"(?<=[.!?])\s+", chunk)
        bullets    = [s.strip() for s in sentences if s.strip()][:6]
        outline.append({
            "slide_number": i,
            "layout":  "BULLETS",
            "title":   first_line,
            "content": bullets,
        })
    outline.append({
        "slide_number": len(outline) + 1,
        "layout":  "CLOSING",
        "title":   "Thank You",
        "content": "Questions & Next Steps",
    })
    return outline


# ─────────────────────────────────────────────────────────────────────────────
# §3.3 — PPTX Assembly (Template Mapping)
# ─────────────────────────────────────────────────────────────────────────────
#
# Layout index mapping per PRD Table 5:
#   TITLE          → 0
#   BULLETS        → 1
#   SECTION_HEADER → 2
#   TWO_COLUMN     → 3
#   CLOSING        → 5
#   QUOTE          → 6
#
# Because python-pptx blank templates don't ship with 6 distinct named layouts,
# we build every slide from the "Blank" layout (index 6 of the default master)
# and draw all visual elements via shapes & text-boxes — giving full visual
# control while faithfully respecting the six layout names and their content
# placeholder specs from the PRD.

def _bg(slide, color: RGBColor):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, fill: RGBColor, line=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    return shape


def _oval(slide, x, y, d, fill: RGBColor):
    shape = slide.shapes.add_shape(9, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _txb(slide, x, y, w, h, text: str, size: float,
         bold=False, italic=False, color: RGBColor = C_WHITE,
         align=PP_ALIGN.LEFT, word_wrap=True) -> None:
    """Add a text-box with a single run."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = "Calibri"

    # PRD §3.3.3 — auto-fit disabled; overflow is logged separately
    tf.auto_size = None


def _check_overflow(slide_title: str, text: str, size_pt: float,
                    box_w_inches: float, box_h_inches: float) -> bool:
    """Rough overflow heuristic — log if content likely overflows."""
    chars_per_line = max(1, int((box_w_inches * 96) / (size_pt * 0.6)))
    lines_needed   = sum(
        max(1, len(line) // chars_per_line + 1) if line else 1
        for line in text.split("\n")
    )
    lines_avail = int((box_h_inches * 96) / (size_pt * 1.4))
    if lines_needed > lines_avail:
        logger.warning(
            f"[OVERFLOW] Slide '{slide_title}': text may exceed placeholder. "
            f"Estimated {lines_needed} lines, box fits ~{lines_avail}."
        )
        return True
    return False


# ── Per-layout renderers ──────────────────────────────────────────────────────

def _render_blank_slide(prs: Presentation) -> any:
    """Add a blank slide (the only reliable way across python-pptx versions)."""
    blank_layout = prs.slide_layouts[6]   # "Blank" is always index 6
    return prs.slides.add_slide(blank_layout)


def render_title(prs: Presentation, refined: dict, slide_data: dict) -> None:
    """PRD Layout TITLE — index 0 semantics: Title (H1) + Subtitle body."""
    slide = _render_blank_slide(prs)
    _bg(slide, C_MIDNIGHT)

    # Full-width left accent stripe
    _rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, C_ACCENT)

    title    = refined.get("title",   slide_data.get("title", "Untitled"))
    subtitle = ""
    bullets  = refined.get("bullets", [])
    if bullets:
        subtitle = bullets[0]
    if not subtitle and isinstance(slide_data.get("content"), str):
        subtitle = slide_data["content"]

    _txb(slide, Inches(0.7), Inches(1.8), Inches(11.8), Inches(1.8),
         title, 44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        _txb(slide, Inches(0.7), Inches(3.8), Inches(11.8), Inches(0.9),
             subtitle, 22, italic=True, color=C_ICE, align=PP_ALIGN.LEFT)

    # Decorative dots
    for i in range(5):
        _oval(slide, Inches(0.7 + i * 0.45), Inches(5.2), Inches(0.14), C_ACCENT)


def render_section_header(prs: Presentation, refined: dict,
                           slide_data: dict) -> None:
    """PRD Layout SECTION_HEADER — index 2: section title + optional number."""
    slide = _render_blank_slide(prs)
    _bg(slide, C_NAVY)

    # Horizontal top bar
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), C_ACCENT)

    title   = refined.get("title", slide_data.get("title", ""))
    bullets = refined.get("bullets", [])
    number  = bullets[0] if bullets else ""

    if number:
        _txb(slide, Inches(0.6), Inches(2.0), Inches(2.5), Inches(1.5),
             number, 72, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
        _txb(slide, Inches(3.2), Inches(2.6), Inches(9.6), Inches(1.2),
             title, 34, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    else:
        _txb(slide, Inches(0.6), Inches(2.5), Inches(12.2), Inches(1.5),
             title, 38, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


def render_bullets(prs: Presentation, refined: dict, slide_data: dict,
                   flagged: list) -> None:
    """
    PRD Layout BULLETS — index 1.
    Title + up to 6 bullet points (PRD §3.3.3).
    Title: 28–36pt bold. Bullets: 18–22pt, 1.2× line spacing (PRD §3.3.3).
    """
    slide  = _render_blank_slide(prs)
    _bg(slide, C_WHITE)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), C_NAVY)

    title   = refined.get("title", slide_data.get("title", ""))
    bullets = refined.get("bullets", [])[:6]   # PRD max 6

    _txb(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.0),
         title, 32, bold=True, color=C_WHITE)

    overflow = _check_overflow(title,
                               "\n".join(bullets), 20, 11.6, 5.6)
    if overflow:
        flagged.append({"slide": title, "reason": "Bullet text may overflow"})

    y = Inches(1.65)
    for bullet in bullets:
        # Accent dot
        _oval(slide, Inches(0.45), y + Inches(0.13), Inches(0.2), C_ACCENT)
        # Bullet text (PRD §3.3.3: 18–22pt, left-aligned)
        _txb(slide, Inches(0.85), y, Inches(11.98), Inches(0.82),
             bullet, 20, color=C_DARK_TXT)
        y += Inches(0.88)


def render_two_column(prs: Presentation, refined: dict,
                      slide_data: dict, flagged: list) -> None:
    """
    PRD Layout TWO_COLUMN — index 3.
    Title + Left column + Right column. Column headers bold (PRD §3.3.3).
    Text: 16–20pt (PRD §3.3.3).
    """
    slide = _render_blank_slide(prs)
    _bg(slide, C_WHITE)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), C_MIDNIGHT)

    title = refined.get("title", slide_data.get("title", ""))
    _txb(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.0),
         title, 32, bold=True, color=C_WHITE)

    # Column backgrounds
    _rect(slide, Inches(0.4),   Inches(1.55), Inches(5.9), Inches(5.65), C_ICE)
    _rect(slide, Inches(7.03),  Inches(1.55), Inches(5.9), Inches(5.65), C_ICE)

    # Column accent tops
    _rect(slide, Inches(0.4),  Inches(1.55), Inches(5.9), Inches(0.12), C_TEAL)
    _rect(slide, Inches(7.03), Inches(1.55), Inches(5.9), Inches(0.12), C_TEAL)

    left_raw  = slide_data.get("content", {})
    left_text  = refined.get("left",  str(left_raw.get("left",  "")) if isinstance(left_raw, dict) else "")
    right_text = refined.get("right", str(left_raw.get("right", "")) if isinstance(left_raw, dict) else "")

    # Default column headers from slide content keys (PRD §3.3.3 — bold headers)
    orig_content = slide_data.get("content", {})
    left_hdr  = "Overview"   if not isinstance(orig_content, dict) else "Left"
    right_hdr = "Details"    if not isinstance(orig_content, dict) else "Right"

    for x, hdr, body in [
        (Inches(0.5),  left_hdr,  left_text),
        (Inches(7.13), right_hdr, right_text),
    ]:
        _txb(slide, x, Inches(1.72), Inches(5.7), Inches(0.55),
             hdr, 18, bold=True, color=C_NAVY)
        _txb(slide, x, Inches(2.35), Inches(5.7), Inches(4.7),
             body, 17, color=C_DARK_TXT)
        overflow = _check_overflow(title, body, 17, 5.7, 4.7)
        if overflow:
            flagged.append({"slide": title,
                            "reason": f"Two-column '{hdr}' text may overflow"})


def render_quote(prs: Presentation, refined: dict, slide_data: dict) -> None:
    """
    PRD Layout QUOTE — index 6.
    Quote text: 24–28pt, italic, centered (PRD §3.3.3).
    Optional attribution.
    """
    slide = _render_blank_slide(prs)
    _bg(slide, C_MIDNIGHT)

    # Large decorative quotation mark
    _txb(slide, Inches(0.4), Inches(0.3), Inches(2.5), Inches(2.0),
         "\u201C", 120, bold=True, color=C_TEAL, align=PP_ALIGN.LEFT)

    bullets = refined.get("bullets", [])
    quote_text = bullets[0] if bullets else slide_data.get("content", "")
    if isinstance(quote_text, list):
        quote_text = " ".join(quote_text)

    attribution = bullets[1] if len(bullets) > 1 else ""

    _txb(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.2),
         str(quote_text), 26, italic=True, color=C_WHITE,
         align=PP_ALIGN.CENTER)

    if attribution:
        _txb(slide, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.7),
             f"— {attribution}", 16, italic=True, color=C_ACCENT,
             align=PP_ALIGN.CENTER)


def render_closing(prs: Presentation, refined: dict, slide_data: dict) -> None:
    """PRD Layout CLOSING — index 5: closing headline + optional body."""
    slide = _render_blank_slide(prs)
    _bg(slide, C_MIDNIGHT)

    _rect(slide, Inches(0), Inches(0),  SLIDE_W, Inches(0.08), C_ACCENT)
    _rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), C_ACCENT)

    title   = refined.get("title", slide_data.get("title", "Thank You"))
    bullets = refined.get("bullets", [])
    body    = bullets[0] if bullets else (
        slide_data.get("content", "") if isinstance(slide_data.get("content"), str)
        else ""
    )

    _txb(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.8),
         title, 48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    if body:
        _txb(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(1.0),
             body, 22, italic=True, color=C_ICE, align=PP_ALIGN.CENTER)

    # Decorative dots
    for i in range(7):
        _oval(slide, Inches(3.2 + i * 1.0), Inches(5.5), Inches(0.14), C_ACCENT)


# ── Slide footer / numbering (PRD §3.3.4) ────────────────────────────────────

def _add_slide_number(slide, number: int, total: int) -> None:
    """Add slide number footer (PRD §3.3.4)."""
    _txb(slide, Inches(11.83), Inches(7.1), Inches(1.4), Inches(0.36),
         f"{number} / {total}", 9, color=C_TEAL, align=PP_ALIGN.RIGHT)


# ── PPTX metadata (PRD §3.3.4) ───────────────────────────────────────────────

def _set_pptx_metadata(prs: Presentation, title: str, author: str = "") -> None:
    core = prs.core_properties
    core.title   = title
    core.author  = author or "Gamma-Local v1.0"
    core.created = datetime.now()
    core.modified = datetime.now()


# ── Main PPTX builder ────────────────────────────────────────────────────────

LAYOUT_RENDERERS = {
    "TITLE":          render_title,
    "SECTION_HEADER": render_section_header,
    "BULLETS":        render_bullets,
    "TWO_COLUMN":     render_two_column,
    "QUOTE":          render_quote,
    "CLOSING":        render_closing,
}

def build_pptx(outline: list[dict], refined_slides: list[dict],
               doc_title: str = "Presentation") -> tuple[bytes, list[dict]]:
    """
    PRD §3.3 — Assemble the PPTX from refined slide content.
    Returns (pptx_bytes, flagged_slides_report).
    """
    prs     = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    _set_pptx_metadata(prs, doc_title)

    flagged: list[dict] = []
    total = len(outline)

    for i, (slide_data, refined) in enumerate(zip(outline, refined_slides), start=1):
        layout = slide_data.get("layout", "BULLETS").upper()
        renderer = LAYOUT_RENDERERS.get(layout, render_bullets)

        # Layouts that don't accept flagged list
        if layout in ("BULLETS", "TWO_COLUMN"):
            renderer(prs, refined, slide_data, flagged)
        else:
            renderer(prs, refined, slide_data)

        # Add slide number footer (PRD §3.3.4)
        slide = prs.slides[-1]
        _add_slide_number(slide, i, total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read(), flagged


# ─────────────────────────────────────────────────────────────────────────────
# §4.4 — Temporary file cleanup
# ─────────────────────────────────────────────────────────────────────────────

def cleanup_temp(path: Path) -> None:
    """Securely delete temp files after pipeline completes (PRD §4.4)."""
    if CFG["output"]["cleanup_temp_files"] and path and path.exists():
        try:
            path.unlink()
            logger.debug(f"Cleaned up temp file: {path}")
        except Exception as e:
            logger.warning(f"Could not delete temp file {path}: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# Audit log (PRD §4.4 — optional, names only, never content)
# ─────────────────────────────────────────────────────────────────────────────

def audit_log(filename: str, model: str, slides_generated: int) -> None:
    if not CFG["security"]["enable_audit_log"]:
        return
    log_path = OUTPUT_DIR / "audit.log"
    entry = {
        "timestamp":        datetime.now().isoformat(),
        "document_name":    filename,   # name only, NO content (PRD §4.4)
        "model":            model,
        "slides_generated": slides_generated,
    }
    with open(log_path, "a") as f:
        f.write(json.dumps(entry) + "\n")


# ─────────────────────────────────────────────────────────────────────────────
# Streamlit UI
# ─────────────────────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="Gamma-Local — AI Presentations",
    page_icon="🛡️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Custom CSS for a polished look
st.markdown("""
<style>
    .stApp { background-color: #0f1923; color: #e8f4fd; }
    .main .block-container { padding-top: 2rem; max-width: 1100px; }
    .metric-card {
        background: #1c3a52; border-radius: 10px;
        padding: 1rem 1.5rem; text-align: center;
        border-left: 4px solid #02c39a;
    }
    .flag-card {
        background: #3a2b1c; border-radius: 8px;
        padding: 0.6rem 1rem; margin-bottom: 0.4rem;
        border-left: 3px solid #f39c12;
    }
    h1 { color: #d6eaf8; }
    .stButton>button {
        background-color: #065a82; color: white;
        border: none; font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)


def sidebar_env_check():
    """PRD §4.3 — environment verification in sidebar."""
    st.sidebar.header("🔍 Environment Check")

    # Check 1: binary
    binary_ok = check_ollama_binary()
    st.sidebar.markdown(
        f"{'✅' if binary_ok else '❌'} Ollama binary "
        f"{'found in PATH' if binary_ok else 'NOT FOUND in PATH'}"
    )
    if not binary_ok:
        st.sidebar.error(
            "Ollama is not installed. Install it:\n\n"
            "```\ncurl -fsSL https://ollama.com/install.sh | sh\n```"
        )

    # Check 2 & 3: server + models
    server_ok, available_models = check_ollama_server()
    st.sidebar.markdown(
        f"{'✅' if server_ok else '❌'} Ollama server "
        f"({'running' if server_ok else 'NOT reachable'} at localhost:11434)"
    )
    if not server_ok:
        st.sidebar.error(
            "Start the Ollama server:\n```\nollama serve\n```"
        )
        return None, []

    supported_available = [
        m for m in available_models
        if any(m.startswith(s.split(":")[0])
               for s in CFG["ollama"]["supported_models"])
    ]

    if not supported_available:
        st.sidebar.error(
            "No supported models found. Pull one:\n\n"
            "```\nollama pull llama3:8b\n```\n\n"
            "Or for lower-spec hardware:\n"
            "```\nollama pull phi3:mini\n```"
        )
        return None, available_models

    best = get_best_available_model(supported_available)
    st.sidebar.markdown(f"✅ Models available: `{'`, `'.join(supported_available)}`")
    return best, available_models


def main():
    # ── Header ────────────────────────────────────────────────────────────────
    col_icon, col_title = st.columns([0.08, 0.92])
    with col_icon:
        st.markdown("# 🛡️")
    with col_title:
        st.markdown("# Gamma-Local")
        st.caption(
            "AI-powered presentations · 100% local · zero data egress · no API keys"
        )

    st.divider()

    # ── Sidebar: env check + model selection ──────────────────────────────────
    best_model, available_models = sidebar_env_check()

    with st.sidebar:
        st.divider()
        st.header("⚙️ Options")

        if available_models:
            supported = available_models  # show everything Ollama has
            chosen_model = st.selectbox(
                "Model",
                options=supported,
                index=(supported.index(best_model)
                       if best_model and best_model in supported else 0),
                help="Recommended: llama3:8b · Fallback: phi3:mini",
            )
        else:
            chosen_model = CFG["ollama"]["default_model"]
            st.info("No models detected. Pull one to begin.")

        slide_count_mode = st.selectbox(
            "Slide count",
            ["Auto (8–20)", "Custom"],
            index=0,
        )
        if slide_count_mode == "Custom":
            target_slides = st.slider(
                "Target slides", 5, CFG["generation"]["max_slides"], 12
            )
        else:
            target_slides = None

        pres_title_override = st.text_input(
            "Presentation title (optional)",
            placeholder="Leave blank to auto-detect from document",
        )

        st.divider()
        st.markdown("""
**Privacy guarantee**
- All inference: localhost only
- No API keys required
- No telemetry or cloud calls
- Temp files auto-deleted
        """)

    # ── Main area: upload ─────────────────────────────────────────────────────
    st.subheader("📂 Upload Document")

    input_mode = st.radio(
        "Input method",
        ["File upload", "Paste text"],
        horizontal=True,
    )

    file_bytes  = None
    filename    = None
    pasted_text = None

    if input_mode == "File upload":
        uploaded = st.file_uploader(
            "Upload PRD, report, or research document",
            type=["pdf", "txt", "md", "docx"],
            help=f"Max size: {CFG['ingestion']['max_file_size_mb']} MB · "
                 "Supported: PDF (text-based), TXT, Markdown, DOCX",
        )
        if uploaded:
            file_bytes = uploaded.read()
            filename   = uploaded.name
            st.success(f"✅ Loaded: **{filename}** "
                       f"({len(file_bytes)/1024:.1f} KB)")

    else:  # Paste text (PRD §3.1.1)
        pasted_text = st.text_area(
            "Paste document text",
            height=250,
            placeholder="Paste your PRD, research brief, or report here…",
        )
        if pasted_text:
            file_bytes = pasted_text.encode("utf-8")
            filename   = "pasted_document.txt"

    # ── Generate button ───────────────────────────────────────────────────────
    generate_btn = st.button(
        "🚀 Generate Presentation",
        type="primary",
        use_container_width=True,
        disabled=(not file_bytes or not best_model),
    )

    if not file_bytes:
        st.info("Upload a document or paste text to begin.")
        return

    if not best_model:
        st.error(
            "No supported Ollama model is available. "
            "Run `ollama pull llama3:8b` then restart the app."
        )
        return

    if not generate_btn:
        return

    # ═════════════════════════════════════════════════════════════════════════
    # PIPELINE
    # ═════════════════════════════════════════════════════════════════════════
    run_id   = str(uuid.uuid4())[:8]
    tmp_path = Path(tempfile.gettempdir()) / f"gamma_local_{run_id}.tmp"
    flagged: list[dict] = []

    progress     = st.progress(0)
    status       = st.empty()
    slide_status = st.empty()

    def upd(pct: int, msg: str):
        progress.progress(pct, text=msg)
        status.markdown(f"**{msg}**")

    try:
        # ── Step 1: Validate ────────────────────────────────────────────────
        upd(5, "📋 Validating file…")
        validate_file(file_bytes, filename)

        # ── Step 2: Extract ─────────────────────────────────────────────────
        upd(10, "📄 Extracting text from document…")
        raw_text = extract_text(file_bytes, filename)

        # ── Step 3: Clean ───────────────────────────────────────────────────
        upd(15, "🧹 Cleaning and normalising text…")
        clean = clean_text(raw_text)

        # ── Step 4: Validate content length (PRD §3.1.3) ───────────────────
        validate_content_length(clean)

        # ── Step 5: Stage 1 — Structural Analysis ──────────────────────────
        upd(20, f"🤖 [{chosen_model}] Analysing document structure…")

        # Inject target_slides hint into prompt if custom mode
        stage1_prompt = clean
        if target_slides:
            stage1_prompt = (
                f"[INSTRUCTION: Generate exactly {target_slides} slides]\n\n"
                + clean
            )

        def stage1_progress(msg: str):
            status.markdown(f"**{msg}**")

        outline = stage1_generate_outline(
            stage1_prompt, chosen_model, stage1_progress
        )
        progress.progress(40)

        doc_title = (
            pres_title_override.strip()
            or outline[0].get("title", "Gamma-Local Presentation")
        )

        # ── Outline preview (PRD §5 Step 4) ────────────────────────────────
        with st.expander(f"📋 Slide Outline Preview — {len(outline)} slides",
                         expanded=True):
            for s in outline:
                icon = {
                    "TITLE":          "🎯",
                    "SECTION_HEADER": "📌",
                    "BULLETS":        "📝",
                    "TWO_COLUMN":     "⚖️",
                    "QUOTE":          "💬",
                    "CLOSING":        "🏁",
                }.get(s.get("layout", ""), "🔹")
                st.markdown(
                    f"{icon} **Slide {s.get('slide_number','?')}** "
                    f"[`{s.get('layout','?')}`] &nbsp; {s.get('title','')}"
                )

        # ── Step 6: Stage 2 — Per-slide content refinement ─────────────────
        refined_slides: list[dict] = []
        total = len(outline)
        for i, slide_data in enumerate(outline, start=1):
            pct = 40 + int((i / total) * 45)
            upd(pct, f"✍️  Refining slide {i} of {total}…")
            slide_status.markdown(
                f"**`{slide_data.get('layout','?')}`** — "
                f"*{slide_data.get('title','')}*"
            )
            refined = stage2_refine_slide(slide_data, chosen_model)
            refined_slides.append(refined)

        slide_status.empty()

        # ── Step 7: PPTX Assembly ───────────────────────────────────────────
        upd(88, "🏗️  Assembling PowerPoint file…")
        pptx_bytes, flagged = build_pptx(outline, refined_slides, doc_title)

        # ── Step 8: Save to output dir (PRD §3.3) ──────────────────────────
        safe_title  = re.sub(r"[^\w\s-]", "", doc_title)[:40].strip().replace(" ", "_")
        output_name = f"{safe_title}_{run_id}.pptx"
        output_path = OUTPUT_DIR / output_name
        output_path.write_bytes(pptx_bytes)
        logger.info(f"PPTX written: {output_path} ({len(pptx_bytes)//1024} KB)")

        audit_log(filename, chosen_model, len(outline))

        upd(100, "✅ Presentation ready!")

        # ── Step 9: Cleanup temp (PRD §4.4) ────────────────────────────────
        cleanup_temp(tmp_path)

    except ValueError as e:
        st.error(f"❌ Input error: {e}")
        logger.error(f"Pipeline ValueError: {e}")
        cleanup_temp(tmp_path)
        return
    except RuntimeError as e:
        st.error(f"❌ Processing error: {e}")
        logger.error(f"Pipeline RuntimeError: {e}")
        cleanup_temp(tmp_path)
        return
    except ollama.ResponseError as e:
        st.error(
            f"❌ Ollama error: {e}\n\n"
            f"Make sure the model `{chosen_model}` is pulled:\n"
            f"```\nollama pull {chosen_model}\n```"
        )
        cleanup_temp(tmp_path)
        return
    except Exception as e:
        st.error(f"❌ Unexpected error: {e}")
        logger.exception("Unhandled pipeline error")
        cleanup_temp(tmp_path)
        return

    # ═════════════════════════════════════════════════════════════════════════
    # Results panel (PRD §5 Step 7)
    # ═════════════════════════════════════════════════════════════════════════
    st.divider()
    st.subheader("🎉 Presentation Ready")

    m1, m2, m3 = st.columns(3)
    with m1:
        st.markdown(
            f'<div class="metric-card"><h2>{len(outline)}</h2>'
            f'<p>Slides generated</p></div>', unsafe_allow_html=True
        )
    with m2:
        st.markdown(
            f'<div class="metric-card"><h2>{len(pptx_bytes)//1024} KB</h2>'
            f'<p>File size</p></div>', unsafe_allow_html=True
        )
    with m3:
        st.markdown(
            f'<div class="metric-card"><h2>{len(flagged)}</h2>'
            f'<p>Slides to review</p></div>', unsafe_allow_html=True
        )

    st.markdown("")

    # Download button (PRD §5 Step 7)
    st.download_button(
        label="⬇️  Download Presentation (.pptx)",
        data=pptx_bytes,
        file_name=output_name,
        mime=(
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
        use_container_width=True,
        type="primary",
    )

    # Slide generation summary report (PRD §5 Step 6)
    with st.expander("📊 Generation Summary Report", expanded=bool(flagged)):
        st.markdown(f"**Document:** `{filename}`")
        st.markdown(f"**Model used:** `{chosen_model}`")
        st.markdown(f"**Output saved to:** `{output_path}`")
        st.markdown(f"**Total slides:** {len(outline)}")

        layout_counts: dict[str, int] = {}
        for s in outline:
            layout_counts[s.get("layout", "?")] = (
                layout_counts.get(s.get("layout", "?"), 0) + 1
            )
        st.markdown("**Layout breakdown:**")
        for layout, count in sorted(layout_counts.items()):
            st.markdown(f"- `{layout}`: {count} slide(s)")

        if flagged:
            st.markdown(
                "### ⚠️ Slides flagged for manual review\n"
                "Text overflow was detected on the following slides. "
                "Open in PowerPoint / LibreOffice to adjust:"
            )
            for f in flagged:
                st.markdown(
                    f'<div class="flag-card">📌 <b>{f["slide"]}</b> — '
                    f'{f["reason"]}</div>',
                    unsafe_allow_html=True,
                )
        else:
            st.success("✅ No overflow issues detected.")

    st.caption(
        "💡 Open the .pptx in PowerPoint, LibreOffice Impress, or Keynote "
        "to fine-tune visuals and add custom charts or images."
    )


if __name__ == "__main__":
    main()